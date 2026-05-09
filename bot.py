import os
import io
import logging
import fitz  # PyMuPDF
from docx import Document
import openpyxl
from pptx import Presentation
from telegram import Update
from telegram.ext import Application, CommandHandler, MessageHandler, filters, ContextTypes
import anthropic

logging.basicConfig(level=logging.INFO)

TELEGRAM_TOKEN = os.environ["TELEGRAM_TOKEN"]
ANTHROPIC_API_KEY = os.environ["ANTHROPIC_API_KEY"]
ADMIN_ID = int(os.environ.get("ADMIN_ID", "0"))

client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

BASE_PROMPT = """Siz "MatbaaPaper | KartonWorks" kompaniyasining HR onboarding assistentidasiz. Ismingiz "Rufat Aminovning agenti".

Xodim o'zbek tilida yozsa — o'zbek tilida, rus tilida yozsa — rus tilida javob ber.

=== KOMPANIYA HAQIDA ===
MatbaaPaper (MP) — makulaturadan gofrolash uchun qog'oz va karton ishlab chiqaradi.
KartonWorks (KW) — tayyor gofrokarton va qadoqlash mahsulotlari ishlab chiqaradi.
Korxona Urgench shahrida joylashgan. 100 dan ortiq xodim.

=== KOMPANIYA FALSAFASI ===
Asosiy tamoyil: "Tizim qahramondan muhimroq. Barqarorlik rekorddan muhimroq."
- QCD (mijozlar): Quality, Cost, Delivery
- PIB (xodimlar): Performance, Income, Balance

=== ISH VAQTI VA TARTIB ===
- Dushanba-Shanba: 09:00-18:00
- Tushlik: 13:00-14:00
- Dam olish: Yakshanba
- Dress code: tartibli, ishchan kiyim majburiy

=== TOPSHIRILADIGAN HUJJATLAR ===
1. Pasport (original + nusxa)
2. Diplom yoki attestat
3. Mehnat daftarchasi
4. PINFL
5. Elektron fotosurat (jpg/png — email orqali)
6. Tibbiy ma'lumotnoma

=== HR BO'LIMI ===
- Aminov Rufat Shavkatovich
- Tel: +998 99 560 74 55
- Email: rufataminov32@gmail.com

=== IMTIYOZLAR ===
- Ish haqi: har oyning 8-10 sanalarida
- Bepul ovqatlanish
- Korporativ aloqa
- Yillik tibbiy ko'rik

=== ISH HAQI TIZIMI ===
Daromad 3 qismdan: Maosh (greyd) + Mukofot (KPI 0-9 ball) + Staj ustamasi (3-12%)
Greydlar: G8=9.6mln, G7.5=8.2mln, G6=5.6mln, G5=4.8mln so'm

=== XAVFSIZLIK ===
- Vision Zero — har qanday baxtsiz hodisaning oldini olish mumkin
- SIZ (himoya vositalari) kiyish majburiy

=== QOIDALAR ===
- Yangi xodimni "Siz" deb murojaat qil
- Javoblar qisqa va aniq (3-5 jumla)
- Bilmagan savolda: Aminov Rufat +998 99 560 74 55
- Har javob oxirida boshqa savol borligini so'ra"""

extra_knowledge = []

def get_system_prompt():
    if extra_knowledge:
        extras = "\n\n=== QOSHIMCHA BILIMLAR ===\n" + "\n".join(
            f"{i+1}. {item}" for i, item in enumerate(extra_knowledge)
        )
        return BASE_PROMPT + extras
    return BASE_PROMPT

user_histories = {}

def is_admin(user_id):
    return user_id == ADMIN_ID

def extract_pdf(data: bytes) -> str:
    doc = fitz.open(stream=data, filetype="pdf")
    return "\n".join(page.get_text() for page in doc)

def extract_docx(data: bytes) -> str:
    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def extract_xlsx(data: bytes) -> str:
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f"[{sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(c) for c in row if c is not None)
            if row_text:
                lines.append(row_text)
    return "\n".join(lines)

def extract_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"[Slayd {i}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text)
    return "\n".join(lines)

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_histories[update.effective_user.id] = []
    await update.message.reply_text(
        "👋 Assalomu alaykum! MatbaaPaper | KartonWorks ga xush kelibsiz!\n\n"
        "Men Rufat Aminovning agenti — onboarding jarayonida sizga ko'maklashaman.\n\n"
        "Savolingizni yozing yoki fayl yuboring (PDF, Word, Excel, PowerPoint) 🇺🇿"
    )

async def help_cmd(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "Nima so'rash mumkin:\n\n"
        "• Qanday hujjatlar kerak?\n"
        "• Ish vaqti va tartib\n"
        "• Ish haqi qachon to'lanadi?\n"
        "• Kompaniya haqida\n"
        "• HR kontakt\n\n"
        "Fayl yuborish: PDF, Word, Excel, PowerPoint\n\n"
        "/start — Suhbatni qayta boshlash"
    )

async def addinfo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not is_admin(update.effective_user.id):
        await update.message.reply_text("Siz admin emassiz.")
        return
    if not context.args:
        await update.message.reply_text("Ishlatish: /addinfo <yangi malumot>")
        return
    text = " ".join(context.args)
    extra_knowledge.append(text)
    await update.message.reply_text(f"Qoshildi ({len(extra_knowledge)}-bilim):\n{text}")

async def showinfo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not is_admin(update.effective_user.id):
        await update.message.reply_text("Siz admin emassiz.")
        return
    if not extra_knowledge:
        await update.message.reply_text("Hozircha qoshimcha bilim yoq.")
        return
    lines = "\n".join(f"{i+1}. {item}" for i, item in enumerate(extra_knowledge))
    await update.message.reply_text(f"Qoshimcha bilimlar:\n\n{lines}")

async def delinfo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not is_admin(update.effective_user.id):
        await update.message.reply_text("Siz admin emassiz.")
        return
    if not context.args:
        await update.message.reply_text("Ishlatish: /delinfo <raqam>")
        return
    try:
        index = int(context.args[0]) - 1
        removed = extra_knowledge.pop(index)
        await update.message.reply_text(f"Ochirildi:\n{removed}")
    except (ValueError, IndexError):
        await update.message.reply_text("Notogri raqam. /showinfo bilan royxatni koring.")

async def myid(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(f"Sizning ID: {update.effective_user.id}")

def clean_markdown(text: str) -> str:
    lines = text.split("\n")
    cleaned = []
    for line in lines:
        if line.startswith("### "):
            line = "*" + line[4:] + "*"
        elif line.startswith("## "):
            line = "*" + line[3:] + "*"
        elif line.startswith("# "):
            line = "*" + line[2:] + "*"
        cleaned.append(line)
    return "\n".join(cleaned)

async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    doc = update.message.document
    fname = doc.file_name.lower()

    if not any(fname.endswith(ext) for ext in [".pdf", ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx"]):
        # Rasm document sifatida yuborilganmi?
        if doc.mime_type in ["image/jpeg", "image/png", "image/jpg"]:
            await handle_image(update, context)
        else:
            await update.message.reply_text("Qabul qilinadigan formatlar: PDF, Word, Excel, PowerPoint, JPEG, PNG")
        return

    await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")
    await update.message.reply_text("Fayl o'qilmoqda...")

    try:
        file = await doc.get_file()
        data = bytes(await file.download_as_bytearray())

        if fname.endswith(".pdf"):
            text = extract_pdf(data)
        elif fname.endswith((".doc", ".docx")):
            text = extract_docx(data)
        elif fname.endswith((".xls", ".xlsx")):
            text = extract_xlsx(data)
        elif fname.endswith((".ppt", ".pptx")):
            text = extract_pptx(data)

        if not text.strip():
            await update.message.reply_text("Faylda matn topilmadi.")
            return

        # Faylni suhbatga qo'shamiz
        uid = update.effective_user.id
        if uid not in user_histories:
            user_histories[uid] = []

        caption = update.message.caption or "Ushbu faylni tahlil qil va asosiy ma'lumotlarni tushuntir."
        user_msg = f"[Fayl: {doc.file_name}]\n\n{text[:8000]}\n\nSavol: {caption}"
        user_histories[uid].append({"role": "user", "content": user_msg})

        response = client.messages.create(
            model="claude-opus-4-7",
            max_tokens=1000,
            system=get_system_prompt(),
            messages=user_histories[uid],
        )
        reply = response.content[0].text
        user_histories[uid].append({"role": "assistant", "content": reply})
        await update.message.reply_text(clean_markdown(reply), parse_mode="Markdown")

    except Exception as e:
        logging.error(f"Fayl xatosi: {e}")
        await update.message.reply_text("Faylni o'qishda xatolik. Qayta urinib ko'ring.")

async def handle_image(update: Update, context: ContextTypes.DEFAULT_TYPE):
    photo = update.message.photo[-1] if update.message.photo else None
    doc = update.message.document

    # Rasm document sifatida yuborilgan bo'lsa
    if doc and doc.mime_type in ["image/jpeg", "image/png", "image/jpg"]:
        file = await doc.get_file()
        mime = doc.mime_type
    elif photo:
        file = await photo.get_file()
        mime = "image/jpeg"
    else:
        return

    await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")
    await update.message.reply_text("Rasm tahlil qilinmoqda...")

    try:
        import base64
        data = bytes(await file.download_as_bytearray())
        b64 = base64.standard_b64encode(data).decode("utf-8")

        uid = update.effective_user.id
        if uid not in user_histories:
            user_histories[uid] = []

        caption = update.message.caption or "Ushbu rasmni tahlil qil va nima ko'rinayotganini tushuntir."

        response = client.messages.create(
            model="claude-opus-4-7",
            max_tokens=1000,
            system=get_system_prompt(),
            messages=user_histories[uid] + [
                {
                    "role": "user",
                    "content": [
                        {"type": "image", "source": {"type": "base64", "media_type": mime, "data": b64}},
                        {"type": "text", "text": caption}
                    ]
                }
            ],
        )
        reply = response.content[0].text
        user_histories[uid].append({"role": "assistant", "content": reply})
        await update.message.reply_text(clean_markdown(reply), parse_mode="Markdown")

    except Exception as e:
        logging.error(f"Rasm xatosi: {e}")
        await update.message.reply_text("Rasmni o'qishda xatolik. Qayta urinib ko'ring.")

async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    text = update.message.text

    if uid not in user_histories:
        user_histories[uid] = []

    user_histories[uid].append({"role": "user", "content": text})

    if len(user_histories[uid]) > 20:
        user_histories[uid] = user_histories[uid][-20:]

    await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")

    try:
        response = client.messages.create(
            model="claude-opus-4-7",
            max_tokens=800,
            system=get_system_prompt(),
            messages=user_histories[uid],
        )
        reply = response.content[0].text
        user_histories[uid].append({"role": "assistant", "content": reply})
        await update.message.reply_text(clean_markdown(reply), parse_mode="Markdown")
    except Exception as e:
        logging.error(f"Error: {e}")
        await update.message.reply_text(
            "Xatolik yuz berdi. Qayta urinib koring.\n"
            "Muammo davom etsa: +998 99 560 74 55"
        )

def main():
    app = Application.builder().token(TELEGRAM_TOKEN).build()
    app.add_handler(CommandHandler("start", start))
    app.add_handler(CommandHandler("help", help_cmd))
    app.add_handler(CommandHandler("addinfo", addinfo))
    app.add_handler(CommandHandler("showinfo", showinfo))
    app.add_handler(CommandHandler("delinfo", delinfo))
    app.add_handler(CommandHandler("myid", myid))
    app.add_handler(MessageHandler(filters.Document.ALL, handle_document))
    app.add_handler(MessageHandler(filters.PHOTO, handle_image))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_message))
    logging.info("Bot ishga tushdi...")
    app.run_polling(drop_pending_updates=True)

if __name__ == "__main__":
    main()
